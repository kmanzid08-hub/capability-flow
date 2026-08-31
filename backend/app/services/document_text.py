import csv
import json
import re
import zipfile
from io import BytesIO
from xml.etree import ElementTree

from docx import Document
from legacy_doc import extract_text as extract_legacy_doc_text
from openpyxl import load_workbook
from pptx import Presentation


class UnsupportedAnalysisDocument(ValueError):
    pass


TEXT_EXTRACTABLE_EXTENSIONS = {
    ".doc",
    ".docx",
    ".txt",
    ".text",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".rtf",
    ".json",
    ".jsonl",
    ".xml",
    ".html",
    ".htm",
    ".yaml",
    ".yml",
    ".xlsx",
    ".xlsm",
    ".pptx",
    ".odt",
    ".ods",
    ".odp",
}

GEMINI_NATIVE_EXTENSIONS = {".pdf", ".jpg", ".jpeg", ".png", ".webp", ".gif"}


def is_gemini_native_document(extension: str) -> bool:
    return extension.lower() in GEMINI_NATIVE_EXTENSIONS


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "cp1252", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _strip_markup(text: str) -> str:
    text = re.sub(r"(?is)<(script|style).*?>.*?</\\1>", " ", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    return re.sub(r"\\s+", " ", text)


def _extract_rtf(content: bytes) -> str:
    text = _decode_text(content)
    text = re.sub(r"\\\\par[d]?\\b", "\n", text)
    text = re.sub(r"\\\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\\\[a-zA-Z]+-?\\d* ?", " ", text)
    return text.replace("{", " ").replace("}", " ").replace("\\", " ")


def _extract_delimited(content: bytes, delimiter: str) -> str:
    reader = csv.reader(_decode_text(content).splitlines(), delimiter=delimiter)
    return "\n".join(" | ".join(cell.strip() for cell in row) for row in reader if row)


def _extract_odf(content: bytes) -> str:
    try:
        with zipfile.ZipFile(BytesIO(content)) as archive:
            xml = archive.read("content.xml")
        root = ElementTree.fromstring(xml)
    except (zipfile.BadZipFile, KeyError, ElementTree.ParseError) as exc:
        raise UnsupportedAnalysisDocument("The OpenDocument file could not be read.") from exc
    chunks: list[str] = []
    for element in root.iter():
        if element.text and element.text.strip():
            chunks.append(element.text.strip())
        if element.tail and element.tail.strip():
            chunks.append(element.tail.strip())
    return "\n".join(chunks)


def extract_text(content: bytes, extension: str, max_chars: int) -> str:
    extension = extension.lower()
    text = ""
    try:
        if extension == ".doc":
            text = extract_legacy_doc_text(content).text
        elif extension == ".docx":
            document = Document(BytesIO(content))
            lines = [p.text for p in document.paragraphs if p.text.strip()]
            for table in document.tables:
                for row in table.rows:
                    values = [c.text.strip() for c in row.cells if c.text.strip()]
                    if values:
                        lines.append(" | ".join(values))
            text = "\n".join(lines)
        elif extension in {".txt", ".text", ".md", ".markdown", ".yaml", ".yml"}:
            text = _decode_text(content)
        elif extension == ".csv":
            text = _extract_delimited(content, ",")
        elif extension == ".tsv":
            text = _extract_delimited(content, "\t")
        elif extension == ".rtf":
            text = _extract_rtf(content)
        elif extension in {".json", ".jsonl"}:
            raw = _decode_text(content)
            if extension == ".json":
                try:
                    text = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
                except json.JSONDecodeError:
                    text = raw
            else:
                text = raw
        elif extension in {".xml", ".html", ".htm"}:
            text = _strip_markup(_decode_text(content))
        elif extension in {".xlsx", ".xlsm"}:
            workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
            lines = []
            for sheet in workbook.worksheets:
                lines.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(v) for v in row if v not in (None, "")]
                    if values:
                        lines.append(" | ".join(values))
            text = "\n".join(lines)
            workbook.close()
        elif extension == ".pptx":
            presentation = Presentation(BytesIO(content))
            lines = []
            for number, slide in enumerate(presentation.slides, start=1):
                lines.append(f"Slide {number}")
                for shape in slide.shapes:
                    value = getattr(shape, "text", "")
                    if value and str(value).strip():
                        lines.append(str(value))
            text = "\n".join(lines)
        elif extension in {".odt", ".ods", ".odp"}:
            text = _extract_odf(content)
        elif is_gemini_native_document(extension):
            raise UnsupportedAnalysisDocument(
                "This file type should be analyzed directly by Gemini rather than "
                "through local text extraction."
            )
        else:
            raise UnsupportedAnalysisDocument(
                f"AI extraction is not available for {extension or 'this file type'} yet."
            )
    except UnsupportedAnalysisDocument:
        raise
    except Exception as exc:
        raise UnsupportedAnalysisDocument(
            f"The {extension or 'uploaded'} file could not be read. "
            "Try opening it and saving it again."
        ) from exc

    normalized = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not normalized:
        raise UnsupportedAnalysisDocument("No readable text was found in this document.")
    return normalized[:max_chars]
