import hashlib
import io
import ipaddress
import socket
from dataclasses import dataclass, field
from pathlib import Path
from urllib.parse import urljoin, urlparse

import httpx
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.core.opportunity_config import get_opportunity_intelligence_settings


class SourceIngestionError(ValueError):
    pass


@dataclass(frozen=True)
class IngestedSource:
    text: str
    content_hash: str
    mime_type: str | None = None
    original_bytes: bytes | None = None
    suggested_filename: str | None = None
    metadata: dict[str, object] = field(default_factory=dict)


class OpportunitySourceIngestionService:
    def __init__(self) -> None:
        self.settings = get_opportunity_intelligence_settings()

    def _finish(
        self,
        text: str,
        mime_type: str | None = None,
        *,
        original_bytes: bytes | None = None,
        suggested_filename: str | None = None,
        metadata: dict[str, object] | None = None,
    ) -> IngestedSource:
        cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
        if len(cleaned) < 20:
            raise SourceIngestionError("The source did not contain enough readable text")
        cleaned = cleaned[: self.settings.opportunity_max_source_characters]
        digest = hashlib.sha256(cleaned.encode("utf-8")).hexdigest()
        return IngestedSource(
            text=cleaned,
            content_hash=digest,
            mime_type=mime_type,
            original_bytes=original_bytes,
            suggested_filename=suggested_filename,
            metadata=metadata or {},
        )

    def _validate_public_url(self, url: str) -> None:
        parsed = urlparse(url)
        if parsed.scheme not in {"http", "https"} or not parsed.hostname:
            raise SourceIngestionError("Only public HTTP and HTTPS URLs are supported")
        try:
            port = parsed.port or (443 if parsed.scheme == "https" else 80)
            addresses = {info[4][0] for info in socket.getaddrinfo(parsed.hostname, port)}
        except socket.gaierror as exc:
            raise SourceIngestionError("The opportunity website could not be resolved") from exc
        for address in addresses:
            ip = ipaddress.ip_address(address)
            if (
                ip.is_private
                or ip.is_loopback
                or ip.is_link_local
                or ip.is_multicast
                or ip.is_reserved
                or ip.is_unspecified
            ):
                raise SourceIngestionError("Private or local network URLs are not allowed")

    async def from_url(self, url: str) -> IngestedSource:
        self._validate_public_url(url)
        headers = {"User-Agent": "CapabilityFlow/1.0 (+opportunity-intake)"}
        async with httpx.AsyncClient(
            timeout=self.settings.opportunity_fetch_timeout_seconds,
            follow_redirects=False,
            headers=headers,
        ) as client:
            current_url = url
            response: httpx.Response | None = None
            for _ in range(6):
                self._validate_public_url(current_url)
                response = await client.get(current_url)
                if response.is_redirect:
                    location = response.headers.get("location")
                    if not location:
                        break
                    current_url = urljoin(str(response.url), location)
                    continue
                response.raise_for_status()
                break
            else:
                raise SourceIngestionError("The opportunity URL redirected too many times")

        if response is None:
            raise SourceIngestionError("The opportunity URL could not be fetched")

        content = response.content
        if len(content) > self.settings.opportunity_max_source_bytes:
            raise SourceIngestionError("The remote opportunity source exceeds the size limit")

        content_type = response.headers.get("content-type", "").split(";", 1)[0].lower()
        final_url = str(response.url)
        path_name = Path(urlparse(final_url).path).name
        metadata: dict[str, object] = {"final_url": final_url}

        if content_type == "application/pdf" or path_name.lower().endswith(".pdf"):
            ingested = self.from_bytes(content, path_name or "source.pdf", content_type)
            return IngestedSource(
                text=ingested.text,
                content_hash=ingested.content_hash,
                mime_type=ingested.mime_type,
                original_bytes=content,
                suggested_filename=path_name or "source.pdf",
                metadata=metadata,
            )

        if content_type not in {"", "text/html", "application/xhtml+xml", "text/plain"}:
            raise SourceIngestionError(f"Unsupported remote content type: {content_type}")

        if content_type == "text/plain":
            return self._finish(
                response.text,
                content_type,
                original_bytes=content,
                suggested_filename=path_name or "source.txt",
                metadata=metadata,
            )

        soup = BeautifulSoup(response.text, "html.parser")
        title = soup.title.get_text(" ", strip=True) if soup.title else None
        canonical = soup.find("link", rel=lambda value: value and "canonical" in value)
        if title:
            metadata["page_title"] = title
        if canonical and canonical.get("href"):
            metadata["canonical_url"] = str(canonical.get("href"))
        for tag in soup(["script", "style", "noscript", "svg", "template"]):
            tag.decompose()
        return self._finish(
            soup.get_text("\n"),
            content_type or "text/html",
            original_bytes=content,
            suggested_filename=path_name or "source.html",
            metadata=metadata,
        )

    def from_bytes(
        self,
        content: bytes,
        filename: str,
        mime_type: str | None = None,
    ) -> IngestedSource:
        if len(content) > self.settings.opportunity_max_source_bytes:
            raise SourceIngestionError("The opportunity document exceeds the size limit")
        suffix = Path(filename).suffix.lower()
        if suffix == ".pdf":
            reader = PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif suffix == ".docx":
            document = Document(io.BytesIO(content))
            doc_chunks = [paragraph.text for paragraph in document.paragraphs]
            for table in document.tables:
                for row in table.rows:
                    doc_chunks.append(" | ".join(cell.text for cell in row.cells))
            text = "\n".join(doc_chunks)
        elif suffix in {".xlsx", ".xlsm"}:
            workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            sheet_chunks: list[str] = []
            for sheet in workbook.worksheets:
                sheet_chunks.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value) for value in row if value not in (None, "")]
                    if values:
                        sheet_chunks.append(" | ".join(values))
            text = "\n".join(sheet_chunks)
        elif suffix == ".pptx":
            presentation = Presentation(io.BytesIO(content))
            slide_chunks: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_chunks.append(str(shape.text))
            text = "\n".join(slide_chunks)
        elif suffix in {".txt", ".csv", ".rtf"}:
            text = content.decode("utf-8", errors="replace")
        else:
            raise SourceIngestionError(f"Unsupported opportunity document type: {suffix or 'none'}")
        return self._finish(
            text,
            mime_type,
            original_bytes=content,
            suggested_filename=filename,
        )
